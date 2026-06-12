"""Controlled in-place edits to Office files (docx/xlsx/pptx).

The sanctioned alternative to ad-hoc interpreter one-liners: every mutating
operation archives a pre-image snapshot before touching the file, so the edit
is reversible with `agw restore` no matter what the library does to the file.

Libraries are optional (python-docx, openpyxl, python-pptx); each operation
reports exactly what to install when its library is missing.
"""
from __future__ import annotations

import os

from core import store


class OfficeError(Exception):
    """Operation failed in a way the caller should report verbatim."""


class MissingLibrary(OfficeError):
    def __init__(self, lib: str, pip_name: str, ext: str):
        super().__init__(f"{ext} support needs the '{pip_name}' package "
                         f"(pip install {pip_name})")
        self.lib, self.pip_name = lib, pip_name


def _docx():
    try:
        import docx
        return docx
    except ImportError:
        raise MissingLibrary("docx", "python-docx", ".docx")


def _openpyxl():
    try:
        import openpyxl
        return openpyxl
    except ImportError:
        raise MissingLibrary("openpyxl", "openpyxl", ".xlsx")


def _pptx():
    try:
        import pptx
        return pptx
    except ImportError:
        raise MissingLibrary("pptx", "python-pptx", ".pptx")


def capabilities() -> dict:
    caps = {}
    for key, loader in (("docx", _docx), ("xlsx", _openpyxl), ("pptx", _pptx)):
        try:
            loader()
            caps[key] = True
        except MissingLibrary:
            caps[key] = False
    return caps


def _ext(path: str) -> str:
    return os.path.splitext(path)[1].lower()


def _snapshot(path: str, op: str) -> dict:
    return store.archive_file(path, mode="copy", dedupe=True,
                              reason=f"pre-image before agw office {op}")


def _iter_pptx_text_frames(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                yield slide, shape.text_frame


def _replace_in_paragraph(paragraph, find: str, replace: str) -> int:
    """Replace within one paragraph, surviving matches that span runs.

    Office splits text into runs at arbitrary points (formatting, spellcheck
    history), so naive per-run replace misses matches. When a match spans
    runs we rewrite the affected runs' text; the first run's formatting wins
    for the replacement text — acceptable for v1, and the pre-image snapshot
    covers regressions.
    """
    full = "".join(run.text for run in paragraph.runs)
    if find not in full:
        return 0
    count = full.count(find)
    new_full = full.replace(find, replace)
    # Re-pack the new text into the existing runs, left to right.
    remaining = new_full
    runs = list(paragraph.runs)
    for i, run in enumerate(runs):
        if i == len(runs) - 1:
            run.text = remaining
            remaining = ""
        else:
            keep = len(run.text)
            run.text, remaining = remaining[:keep], remaining[keep:]
    return count


# --- read operations (no snapshot needed) -------------------------------------

def get_text(path: str) -> str:
    ext = _ext(path)
    if ext == ".docx":
        doc = _docx().Document(path)
        parts = [p.text for p in doc.paragraphs]
        for table in doc.tables:
            for row in table.rows:
                parts.append("\t".join(c.text for c in row.cells))
        return "\n".join(parts)
    if ext == ".pptx":
        prs = _pptx().Presentation(path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"--- slide {i} ---")
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text:
                    parts.append(shape.text_frame.text)
        return "\n".join(parts)
    raise OfficeError(f"get-text supports .docx and .pptx (for {ext or 'this file'}, "
                      "use `agw checkout` / `agw convert`)")


def info(path: str) -> dict:
    ext = _ext(path)
    if ext == ".docx":
        doc = _docx().Document(path)
        headings = [p.text for p in doc.paragraphs
                    if p.style.name.startswith("Heading") and p.text.strip()]
        return {"type": "docx", "paragraphs": len(doc.paragraphs),
                "tables": len(doc.tables), "headings": headings}
    if ext in (".xlsx", ".xlsm"):
        wb = _openpyxl().load_workbook(path, read_only=True)
        sheets = {ws.title: {"rows": ws.max_row, "cols": ws.max_column}
                  for ws in wb.worksheets}
        wb.close()
        return {"type": "xlsx", "sheets": sheets}
    if ext == ".pptx":
        prs = _pptx().Presentation(path)
        titles = []
        for slide in prs.slides:
            title = slide.shapes.title
            titles.append(title.text if title is not None else "")
        return {"type": "pptx", "slides": len(prs.slides), "titles": titles}
    raise OfficeError(f"unsupported extension: {ext or 'none'}")


# --- write operations (snapshot first, always) ---------------------------------

def replace_text(path: str, find: str, replace: str) -> dict:
    if not find:
        raise OfficeError("--find must not be empty")
    ext = _ext(path)
    if ext == ".docx":
        docx = _docx()
        doc = docx.Document(path)
        snap = _snapshot(path, "replace-text")
        n = 0
        paragraphs = list(doc.paragraphs)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    paragraphs.extend(cell.paragraphs)
        for p in paragraphs:
            n += _replace_in_paragraph(p, find, replace)
        if n:
            doc.save(path)
        return {"replacements": n, "snapshot_version": snap.get("version")}
    if ext == ".pptx":
        pptx = _pptx()
        prs = pptx.Presentation(path)
        snap = _snapshot(path, "replace-text")
        n = 0
        for _, tf in _iter_pptx_text_frames(prs):
            for p in tf.paragraphs:
                n += _replace_in_paragraph(p, find, replace)
        if n:
            prs.save(path)
        return {"replacements": n, "snapshot_version": snap.get("version")}
    raise OfficeError(f"replace-text supports .docx and .pptx, not {ext or 'this file'}")


def _coerce(value: str, force_text: bool):
    if force_text:
        return value
    if value.startswith("="):
        return value  # formula
    for caster in (int, float):
        try:
            return caster(value)
        except ValueError:
            continue
    if value.lower() in ("true", "false"):
        return value.lower() == "true"
    return value


def set_cell(path: str, sheet: str, cell: str, value: str,
             force_text: bool = False) -> dict:
    openpyxl = _openpyxl()
    wb = openpyxl.load_workbook(path)
    if sheet not in wb.sheetnames:
        raise OfficeError(f"no sheet named {sheet!r} (have: {', '.join(wb.sheetnames)})")
    snap = _snapshot(path, "set-cell")
    ws = wb[sheet]
    old = ws[cell].value
    ws[cell] = _coerce(value, force_text)
    wb.save(path)
    return {"sheet": sheet, "cell": cell, "old": old, "new": ws[cell].value,
            "snapshot_version": snap.get("version")}


def append_rows(path: str, sheet: str, rows: list, force_text: bool = False) -> dict:
    openpyxl = _openpyxl()
    wb = openpyxl.load_workbook(path)
    if sheet not in wb.sheetnames:
        raise OfficeError(f"no sheet named {sheet!r} (have: {', '.join(wb.sheetnames)})")
    snap = _snapshot(path, "append-rows")
    ws = wb[sheet]
    for row in rows:
        ws.append([_coerce(str(v), force_text) if v is not None else None
                   for v in row])
    wb.save(path)
    return {"sheet": sheet, "appended": len(rows), "rows_now": ws.max_row,
            "snapshot_version": snap.get("version")}
